"""Convert PPTX to markdown text. Images are replaced with a placeholder."""
import pathlib
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def pptx_to_markdown(pad: pathlib.Path) -> str:
    prs = Presentation(pad)
    parts = []

    for i, slide in enumerate(prs.slides, 1):
        slide_parts = [f"## Slide {i}"]
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_parts.append("[AFBEELDING VERWIJDERD]")
            elif shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    tekst = para.text.strip()
                    if tekst:
                        slide_parts.append(tekst)
        if len(slide_parts) > 1:
            parts.append("\n".join(slide_parts))

    return "\n\n".join(parts)
